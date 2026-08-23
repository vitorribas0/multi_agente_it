from django.conf import settings
from django.shortcuts import render, get_object_or_404
from django.http import JsonResponse, FileResponse, Http404, StreamingHttpResponse
from django.views.decorators.csrf import csrf_exempt
from django.views.decorators.http import require_GET, require_http_methods
import io
import json
import os
import queue
import re
import threading
from pathlib import Path

from .ai_service import run_agent, request_stop, register_stop, clear_stop, MODEL_OPTIONS
from .models import (
    Conversation, Message, ToolCall, Agent, SessionAgent, AppSettings,
    Knowledge, Playbook,
)
from tools import all_tools


# Limites do upload de tabelas
UPLOAD_MAX_BYTES = 1024 * 1024 * 1024    # 1 GB de arquivo
UPLOAD_MAX_ROWS = 2_000_000               # nº máx. de linhas processadas
PREVIEW_ROWS_FOR_LLM = 3                  # quantas linhas vão pro modelo
PREVIEW_ROWS_FOR_UI = 100                 # tamanho da página inicial na UI

TABLE_EXTS = {".csv", ".xlsx", ".xls", ".parquet"}
DOCUMENT_EXTS = {".pdf", ".docx", ".doc", ".pptx", ".ppt",
                 ".html", ".htm", ".md", ".txt",
                 ".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp"}
BATCH_PDF_EXTS = {".pdf", ".txt"}  # extensões aceitas no upload em lote
BATCH_MAX_FILES = 200              # nº máx. de arquivos por upload em lote


# ── Páginas ──────────────────────────────────────────────────────────

def index(request):
    return render(request, "auditor/index.html")


def manual(request):
    return render(request, "auditor/manual.html")


def settings_page(request):
    return render(request, "auditor/settings.html")


# ── Helpers de serialização ──────────────────────────────────────────

def _conversation_summary(c: Conversation) -> dict:
    return {
        "id": c.id,
        "title": c.title,
        "agent_slug": c.agent.slug if c.agent else None,
        "awaiting_human_input": c.awaiting_human_input,
        "has_session_agent": hasattr(c, "session_agent"),
        "playbook_id": c.playbook_id,
        "playbook_name": c.playbook.name if c.playbook_id else None,
        "updated_at": c.updated_at.isoformat(),
    }


def _session_agent_payload(sa: SessionAgent) -> dict:
    return {
        "name": sa.name,
        "icon": sa.icon,
        "system_prompt": sa.system_prompt,
        "model": sa.model,
        "temperature": sa.temperature,
        "tools_enabled": sa.tools_enabled or [],
        "guardrails": sa.guardrails,
        "documents": _documents_summary(sa.documents or []),
    }


def _documents_summary(documents: list) -> list:
    """Versão enxuta dos documentos para o frontend — sem o markdown completo
    (que pode ser grande). O conteúdo só trafega no upload e fica no banco."""
    out = []
    for d in documents:
        if not isinstance(d, dict):
            continue
        out.append({
            "filename": d.get("filename", "documento"),
            "char_count": d.get("char_count", 0),
            "page_count": d.get("page_count"),
        })
    return out


def _message_payload(m: Message) -> dict:
    return {
        "id": m.id,
        "role": m.role,
        "content": m.content,
        "attachment": m.attachment or None,
        "attachments": m.attachments or [],
        "tool_calls": [
            {
                "tool": tc.tool_name,
                "args": tc.args,
                "result": tc.result,
                "error": tc.error,
                "duration_ms": tc.duration_ms,
                "nested_tool_calls": tc.nested_tool_calls or [],
            }
            for tc in m.tool_calls.all()
        ],
    }


def _artifacts_ref_block(attachments) -> str:
    """Resume os artefatos de UMA mensagem como texto para reinjetar no histórico.

    O modelo salva arquivos (HTML/PDF/CSV/gráfico) via tools, mas o histórico
    reenviado a cada turno só levava `content` — as referências (filename +
    download_url) ficavam presas em `Message.attachments` e o modelo "esquecia"
    o que já havia produzido. Aqui reconstruímos uma linha por artefato para que
    ele lembre, em turnos seguintes, dos arquivos que gerou e onde estão.

    Não injetamos payloads pesados (ex.: base64 de gráfico) — apenas metadados
    leves que permitem o modelo referenciar o artefato.
    """
    if not attachments:
        return ""
    lines = []
    for a in attachments:
        if not isinstance(a, dict):
            continue
        kind = a.get("kind") or "artefato"
        nome = a.get("filename") or a.get("titulo") or "(sem nome)"
        url = a.get("download_url")
        if url:
            extra = f" ({a.get('formato')}, {a.get('size_kb')} KB)" if a.get("formato") else ""
            lines.append(f"- {kind}: \"{nome}\"{extra} → {url}")
        elif kind == "table":
            rows, cols = a.get("rows"), a.get("columns")
            n_cols = len(cols) if isinstance(cols, list) else cols
            lines.append(f"- table: \"{nome}\" ({rows} linhas × {n_cols} colunas)")
        elif kind == "chart":
            lines.append(f"- chart: \"{nome}\" (gráfico {a.get('chart_type') or ''})".strip())
        elif kind == "mermaid":
            lines.append(f"- diagrama: \"{nome}\"")
        else:
            lines.append(f"- {kind}: \"{nome}\"")
    if not lines:
        return ""
    return "[Artefatos gerados neste turno — disponíveis para referência:\n" + "\n".join(lines) + "]"


def _build_history(conv) -> list[dict]:
    """Monta o histórico enviado ao modelo, reanexando referências de artefatos.

    Sem isso o modelo não enxerga, em turnos seguintes, os arquivos que gerou
    (ver `_artifacts_ref_block`). Usado tanto por `chat_message` quanto por
    `chat_stream` — antes o mapeamento estava duplicado nos dois.
    """
    history = []
    for m in conv.messages.all():
        content = m.content or ""
        if m.role == "assistant":
            refs = _artifacts_ref_block(m.attachments)
            if refs:
                content = f"{content}\n\n{refs}" if content else refs
        history.append({"role": m.role, "content": content})
    return history


# ── Conversas ────────────────────────────────────────────────────────

@require_GET
def conversation_list(request):
    # select_related traz agent e session_agent na MESMA query.
    # Sem isso, cada _conversation_summary dispara 2 queries por conversa
    # (c.agent e o reverse OneToOne session_agent), deixando o histórico
    # cada vez mais lento conforme o número de chats cresce.
    convs = Conversation.objects.select_related("agent", "session_agent", "playbook")
    return JsonResponse({"conversations": [_conversation_summary(c) for c in convs]})


@require_GET
def conversation_detail(request, conv_id):
    # prefetch_related("tool_calls") evita o N+1 ao montar as mensagens:
    # sem ele, cada _message_payload dispara uma query para m.tool_calls,
    # deixando a abertura do chat lenta quando há muitas mensagens.
    conv = get_object_or_404(
        Conversation.objects.select_related("agent", "session_agent", "playbook"),
        id=conv_id,
    )
    messages = conv.messages.prefetch_related("tool_calls")
    return JsonResponse({
        "id": conv.id,
        "title": conv.title,
        "agent_slug": conv.agent.slug if conv.agent else None,
        "awaiting_human_input": conv.awaiting_human_input,
        "has_session_agent": hasattr(conv, "session_agent"),
        "playbook_id": conv.playbook_id,
        "playbook_name": conv.playbook.name if conv.playbook_id else None,
        "messages": [_message_payload(m) for m in messages],
    })


@csrf_exempt
@require_http_methods(["PATCH", "POST"])
def conversation_rename(request, conv_id):
    """Renomeia uma conversa (atualiza o título)."""
    conv = get_object_or_404(Conversation, id=conv_id)
    try:
        data = json.loads(request.body or "{}")
        title = str(data.get("title") or "").strip()
        if not title:
            return JsonResponse(
                {"status": "error", "message": "Título vazio."}, status=400
            )
        conv.title = title[:120]
        conv.save(update_fields=["title", "updated_at"])
        return JsonResponse({"status": "success", "conversation": _conversation_summary(conv)})
    except Exception as e:
        return JsonResponse({"status": "error", "message": str(e)}, status=400)


@csrf_exempt
@require_http_methods(["DELETE"])
def conversation_delete(request, conv_id):
    conv = get_object_or_404(Conversation, id=conv_id)
    conv.delete()
    return JsonResponse({"status": "success"})


# ── Knowledge Bases (RAG) ────────────────────────────────────────────

import time as _time

# Cache simples em memória da listagem de KBs (a chamada IARA não é barata).
_KB_CACHE = {"ts": 0.0, "items": []}
_KB_CACHE_TTL_S = 300  # 5 min
KB_MAX = 10            # máximo de KBs ativas por conversa


def _list_kbs(force_refresh: bool = False) -> list:
    """Lista as Knowledge Bases do projeto IARA, com cache em memória."""
    now = _time.time()
    if not force_refresh and _KB_CACHE["items"] and (now - _KB_CACHE["ts"] < _KB_CACHE_TTL_S):
        return _KB_CACHE["items"]

    from uuid import uuid4
    from iaragenai import IaraGenAI

    client = IaraGenAI(
        client_id=os.getenv("IARA_CLIENT_ID"),
        client_secret=os.getenv("IARA_CLIENT_SECRET"),
        environment=os.getenv("IARA_ENVIRONMENT", "homol"),
        access_token=os.getenv("IARA_ACCESS_TOKEN") or os.getenv("ACCESS_TOKEN"),
        correlation_id=str(uuid4()),
    )
    raw = client.knowledge_base.list() or []
    items = [
        {
            "id": getattr(kb, "knowledge_base_id", None),
            "name": getattr(kb, "name", "") or "",
            "description": getattr(kb, "description", "") or "",
        }
        for kb in raw
        if getattr(kb, "knowledge_base_id", None)
    ]
    _KB_CACHE["items"] = items
    _KB_CACHE["ts"] = now
    return items


@require_GET
def kbs_list(request):
    """Lista as Knowledge Bases disponíveis no projeto IARA."""
    try:
        force = request.GET.get("refresh") == "1"
        return JsonResponse({"status": "success", "kbs": _list_kbs(force_refresh=force)})
    except Exception as e:
        return JsonResponse({"status": "error", "message": str(e)}, status=500)


@require_GET
def conversation_kbs(request, conv_id):
    """Retorna as KBs ativas salvas na conversa."""
    conv = get_object_or_404(Conversation, id=conv_id)
    active = (conv.state or {}).get("active_kbs") or []
    return JsonResponse({"status": "success", "active_kbs": active})


def _normalize_active_kbs(incoming) -> list:
    """Normaliza a lista de KBs vinda do cliente (limita a KB_MAX itens válidos)."""
    if not isinstance(incoming, list):
        return []
    return [
        {
            "id": k.get("id"),
            "name": k.get("name", ""),
            "description": k.get("description", ""),
        }
        for k in incoming
        if isinstance(k, dict) and k.get("id")
    ][:KB_MAX]


def _apply_active_kbs(conv, incoming) -> None:
    """Persiste active_kbs em conv.state. None = não mexe na seleção atual."""
    if incoming is None:
        return
    state = dict(conv.state or {})
    state["active_kbs"] = _normalize_active_kbs(incoming)
    conv.state = state
    conv.save(update_fields=["state", "updated_at"])


@csrf_exempt
@require_http_methods(["POST", "PUT"])
def conversation_kbs_save(request, conv_id):
    """Salva a seleção de KBs ativas na conversa (Conversation.state)."""
    conv = get_object_or_404(Conversation, id=conv_id)
    try:
        data = json.loads(request.body or "{}")
        active = _normalize_active_kbs(data.get("active_kbs") or [])
        state = dict(conv.state or {})
        state["active_kbs"] = active
        conv.state = state
        conv.save(update_fields=["state", "updated_at"])
        return JsonResponse({"status": "success", "active_kbs": active})
    except Exception as e:
        return JsonResponse({"status": "error", "message": str(e)}, status=400)


# ── Conhecimentos (prompts de especialista cadastrados na tela) ──────

KNOWLEDGE_MAX = 10  # máximo de conhecimentos ativos por conversa


def _knowledge_payload(k) -> dict:
    return {
        "id": k.id,
        "name": k.name,
        "description": k.description,
        "icon": k.icon,
        "prompt": k.prompt,
    }


@require_GET
def knowledge_list(request):
    """Lista todos os conhecimentos cadastrados."""
    items = [_knowledge_payload(k) for k in Knowledge.objects.all()]
    return JsonResponse({"status": "success", "knowledge": items})


@csrf_exempt
@require_http_methods(["POST"])
def knowledge_create(request):
    """Cria um novo conhecimento."""
    try:
        data = json.loads(request.body or "{}")
        name = (data.get("name") or "").strip()
        prompt = (data.get("prompt") or "").strip()
        if not name:
            return JsonResponse({"status": "error", "message": "Nome é obrigatório."}, status=400)
        if not prompt:
            return JsonResponse({"status": "error", "message": "Prompt é obrigatório."}, status=400)
        k = Knowledge.objects.create(
            name=name[:120],
            description=(data.get("description") or "").strip()[:240],
            icon=(data.get("icon") or "📚").strip()[:8] or "📚",
            prompt=prompt,
        )
        return JsonResponse({"status": "success", "knowledge": _knowledge_payload(k)})
    except Exception as e:
        return JsonResponse({"status": "error", "message": str(e)}, status=400)


@csrf_exempt
@require_http_methods(["POST"])
def knowledge_update(request, know_id):
    """Atualiza um conhecimento existente."""
    k = get_object_or_404(Knowledge, id=know_id)
    try:
        data = json.loads(request.body or "{}")
        if "name" in data:
            name = (data.get("name") or "").strip()
            if not name:
                return JsonResponse({"status": "error", "message": "Nome é obrigatório."}, status=400)
            k.name = name[:120]
        if "description" in data:
            k.description = (data.get("description") or "").strip()[:240]
        if "icon" in data:
            k.icon = (data.get("icon") or "📚").strip()[:8] or "📚"
        if "prompt" in data:
            prompt = (data.get("prompt") or "").strip()
            if not prompt:
                return JsonResponse({"status": "error", "message": "Prompt é obrigatório."}, status=400)
            k.prompt = prompt
        k.save()
        return JsonResponse({"status": "success", "knowledge": _knowledge_payload(k)})
    except Exception as e:
        return JsonResponse({"status": "error", "message": str(e)}, status=400)


@csrf_exempt
@require_http_methods(["POST", "DELETE"])
def knowledge_delete(request, know_id):
    """Exclui um conhecimento."""
    k = get_object_or_404(Knowledge, id=know_id)
    try:
        k.delete()
        return JsonResponse({"status": "success"})
    except Exception as e:
        return JsonResponse({"status": "error", "message": str(e)}, status=400)


@require_GET
def conversation_knowledge(request, conv_id):
    """Retorna os conhecimentos ativos salvos na conversa."""
    conv = get_object_or_404(Conversation, id=conv_id)
    active = (conv.state or {}).get("active_knowledge") or []
    return JsonResponse({"status": "success", "active_knowledge": active})


def _normalize_active_knowledge(incoming) -> list:
    """Normaliza a lista de conhecimentos ativos vinda do cliente.

    Só guardamos o id na conversa; o conteúdo do prompt é sempre lido do
    banco na hora de executar (assim edições no conhecimento valem na hora e
    não incham o Conversation.state). Limita a KNOWLEDGE_MAX itens válidos.
    """
    if not isinstance(incoming, list):
        return []
    seen, out = set(), []
    for item in incoming:
        kid = item.get("id") if isinstance(item, dict) else item
        try:
            kid = int(kid)
        except (TypeError, ValueError):
            continue
        if kid in seen:
            continue
        seen.add(kid)
        out.append({"id": kid})
    return out[:KNOWLEDGE_MAX]


def _apply_active_knowledge(conv, incoming) -> None:
    """Persiste active_knowledge em conv.state. None = não mexe na seleção."""
    if incoming is None:
        return
    state = dict(conv.state or {})
    state["active_knowledge"] = _normalize_active_knowledge(incoming)
    conv.state = state
    conv.save(update_fields=["state", "updated_at"])


def _resolve_active_knowledge(conv) -> list:
    """Lê os ids ativos da conversa e busca os conhecimentos no banco.

    Retorna [{id, name, description, prompt}] na ordem salva, ignorando ids
    que já não existem (conhecimento excluído). O conteúdo vem sempre do banco
    para que edições valham imediatamente.
    """
    active = (conv.state or {}).get("active_knowledge") or []
    ids = [a.get("id") for a in active if isinstance(a, dict) and a.get("id")]
    if not ids:
        return []
    by_id = {k.id: k for k in Knowledge.objects.filter(id__in=ids)}
    resolved = []
    for kid in ids:
        k = by_id.get(kid)
        if k is not None:
            resolved.append({
                "id": k.id,
                "name": k.name,
                "description": k.description,
                "prompt": k.prompt,
            })
    return resolved


@csrf_exempt
@require_http_methods(["POST", "PUT"])
def conversation_knowledge_save(request, conv_id):
    """Salva a seleção de conhecimentos ativos na conversa (state)."""
    conv = get_object_or_404(Conversation, id=conv_id)
    try:
        data = json.loads(request.body or "{}")
        active = _normalize_active_knowledge(data.get("active_knowledge") or [])
        state = dict(conv.state or {})
        state["active_knowledge"] = active
        conv.state = state
        conv.save(update_fields=["state", "updated_at"])
        return JsonResponse({"status": "success", "active_knowledge": active})
    except Exception as e:
        return JsonResponse({"status": "error", "message": str(e)}, status=400)


# ── Chat ─────────────────────────────────────────────────────────────

def _resolve_conversation(conv_id, agent_slug, text):
    """Recupera ou cria a conversa. Retorna (conv, erro_str)."""
    conv: Conversation | None = None
    if conv_id:
        conv = Conversation.objects.filter(id=conv_id).first()
    if conv is None:
        agent = None
        if agent_slug:
            agent = Agent.objects.filter(slug=agent_slug).first()
        if agent is None:
            agent = Agent.objects.filter(is_default=True).first() or Agent.objects.first()
        if agent is None:
            return None, "Nenhum agente cadastrado. Acesse Configurações."
        title = text[:60] + ("…" if len(text) > 60 else "")
        conv = Conversation.objects.create(title=title, agent=agent)
    if conv.agent is None:
        return None, "Conversa sem agente associado."
    return conv, None


def _persist_turn(conv, session_state, result):
    """Salva a resposta do assistente + tool calls e atualiza a conversa.

    Remove os campos de controle internos (prefixo '__') do estado antes de
    persistir — eles incluem o callback de progresso (não serializável) e o
    histórico, que não devem ir para o JSONField.
    """
    pending_attachments = session_state.pop("__pending_attachments", None) or []
    for k in [k for k in session_state if k.startswith("__")]:
        session_state.pop(k, None)

    assistant_msg = Message.objects.create(
        conversation=conv,
        role="assistant",
        content=result.answer or (result.human_question if result.awaiting_human else ""),
        attachments=pending_attachments,
    )
    for tc in result.tool_calls:
        ToolCall.objects.create(
            message=assistant_msg,
            tool_name=tc["tool"],
            args=tc.get("args") or {},
            result=tc.get("result") or "",
            error=tc.get("error") or "",
            duration_ms=tc.get("duration_ms") or 0,
            nested_tool_calls=tc.get("nested_tool_calls") or [],
        )

    if result.state_changed:
        conv.state = session_state
    conv.awaiting_human_input = result.awaiting_human
    conv.save()
    return assistant_msg


@csrf_exempt
def chat_message(request):
    if request.method != "POST":
        return JsonResponse({"status": "error", "message": "Method not allowed"}, status=405)

    try:
        data = json.loads(request.body)
        text = (data.get("message") or "").strip()
        conv_id = data.get("conversation_id")
        agent_slug = data.get("agent_slug")  # opcional ao criar nova conversa

        if not text:
            return JsonResponse({"status": "error", "message": "Mensagem vazia."}, status=400)

        conv, err = _resolve_conversation(conv_id, agent_slug, text)
        if err:
            return JsonResponse({"status": "error", "message": err}, status=400)

        _apply_active_kbs(conv, data.get("active_kbs"))
        _apply_active_knowledge(conv, data.get("active_knowledge"))
        _apply_playbook(conv, data.get("playbook_id", _UNSET))

        history = _build_history(conv)
        Message.objects.create(conversation=conv, role="user", content=text)

        session_state = dict(conv.state or {})
        session_state["__conversation_id"] = conv.id
        session_state["__active_knowledge"] = _resolve_active_knowledge(conv)
        pb_root, pb_extras = _playbook_runtime(conv)
        session_state.update(pb_extras)
        result = run_agent(
            agent=pb_root if pb_root is not None else conv.agent,
            user_message=text,
            history=history,
            session=session_state,
        )
        assistant_msg = _persist_turn(conv, session_state, result)

        return JsonResponse({
            "status": "success",
            "conversation_id": conv.id,
            "conversation_title": conv.title,
            "agent_slug": conv.agent.slug,
            "awaiting_human_input": result.awaiting_human,
            "human_question": result.human_question,
            "reply": _message_payload(assistant_msg),
        })
    except Exception as e:
        import traceback
        return JsonResponse(
            {"status": "error", "message": str(e), "traceback": traceback.format_exc()},
            status=500,
        )


@csrf_exempt
def chat_stream(request):
    """Versão streaming (SSE) do chat: empurra eventos de progresso ao vivo.

    O agente roda numa thread separada e publica eventos numa fila; o
    generator drena a fila e emite frames SSE ('data: {...}\\n\\n'). O último
    evento ('done' ou 'error') carrega o payload final igual ao chat_message.
    """
    if request.method != "POST":
        return JsonResponse({"status": "error", "message": "Method not allowed"}, status=405)

    data = json.loads(request.body)
    text = (data.get("message") or "").strip()
    conv_id = data.get("conversation_id")
    agent_slug = data.get("agent_slug")

    if not text:
        return JsonResponse({"status": "error", "message": "Mensagem vazia."}, status=400)

    conv, err = _resolve_conversation(conv_id, agent_slug, text)
    if err:
        return JsonResponse({"status": "error", "message": err}, status=400)

    _apply_active_kbs(conv, data.get("active_kbs"))
    _apply_active_knowledge(conv, data.get("active_knowledge"))
    _apply_playbook(conv, data.get("playbook_id", _UNSET))

    history = _build_history(conv)
    Message.objects.create(conversation=conv, role="user", content=text)

    # Resolve os conhecimentos ativos aqui (thread principal, ORM disponível);
    # a thread worker apenas consome a lista já pronta.
    active_knowledge = _resolve_active_knowledge(conv)
    # Idem para o runtime do playbook: monta o nó root + extras de sessão na
    # thread principal (a worker só consome). Reload garante que uma vinculação
    # recém-aplicada acima seja enxergada.
    conv.refresh_from_db(fields=["playbook"])
    pb_root, pb_extras = _playbook_runtime(conv)

    events: "queue.Queue[dict]" = queue.Queue()
    SENTINEL = {"__sentinel__": True}

    # Registra o canal de stop ANTES de iniciar a thread, para que um clique
    # imediato no botão de parar já encontre o event registrado.
    stop_event = register_stop(conv.id)

    def worker():
        session_state = dict(conv.state or {})
        session_state["__conversation_id"] = conv.id
        session_state["__stop_event"] = stop_event
        session_state["__active_knowledge"] = active_knowledge
        session_state.update(pb_extras)

        def progress(payload: dict) -> None:
            events.put({"type": "progress", **payload})

        try:
            result = run_agent(
                agent=pb_root if pb_root is not None else conv.agent,
                user_message=text,
                history=history,
                session=session_state,
                progress=progress,
            )
            assistant_msg = _persist_turn(conv, session_state, result)
            events.put({
                "type": "done",
                "payload": {
                    "status": "success",
                    "conversation_id": conv.id,
                    "conversation_title": conv.title,
                    "agent_slug": conv.agent.slug,
                    "awaiting_human_input": result.awaiting_human,
                    "human_question": result.human_question,
                    "stopped": result.stopped,
                    "reply": _message_payload(assistant_msg),
                },
            })
        except Exception as e:
            import traceback
            traceback.print_exc()
            events.put({"type": "error", "message": str(e)})
        finally:
            clear_stop(conv.id)
            # Thread própria => conexão de DB própria; fecha pra não vazar.
            from django.db import connection
            connection.close()
            events.put(SENTINEL)

    def stream():
        # Comentário inicial desabilita o buffering de alguns proxies.
        yield ": stream start\n\n"
        threading.Thread(target=worker, daemon=True).start()
        while True:
            try:
                evt = events.get(timeout=15)
            except queue.Empty:
                # Heartbeat para manter a conexão viva atrás de proxies.
                yield ": keep-alive\n\n"
                continue
            if evt is SENTINEL:
                break
            yield f"data: {json.dumps(evt, ensure_ascii=False)}\n\n"

    response = StreamingHttpResponse(stream(), content_type="text/event-stream")
    response["Cache-Control"] = "no-cache"
    response["X-Accel-Buffering"] = "no"  # desliga buffering do nginx, se houver
    response["X-Conversation-Id"] = str(conv.id)
    return response


@csrf_exempt
@require_http_methods(["POST"])
def chat_stop(request, conv_id):
    """Interrompe a geração em andamento de uma conversa.

    Seta o event de stop registrado pelo turno corrente; o loop do agente o
    detecta no próximo passo e encerra preservando os resultados parciais.
    """
    legacy_stopped = request_stop(conv_id)
    # Import local evita ciclo de módulos: codex_views reutiliza helpers desta
    # view, enquanto este endpoint unifica o botão Parar dos dois motores.
    from .codex_views import request_codex_stop

    codex_stopped = request_codex_stop(conv_id)
    return JsonResponse({
        "status": "success",
        "stopped": legacy_stopped or codex_stopped,
        "legacy_stopped": legacy_stopped,
        "codex_stopped": codex_stopped,
    })


# ── Configuração: agentes, modelos e tools ──────────────────────────

@require_GET
def config_overview(request):
    """Retorna agentes, modelos e tools disponíveis."""
    agents_data = [
        {
            "id": a.id,
            "slug": a.slug,
            "name": a.name,
            "description": a.description,
            "icon": a.icon,
            "system_prompt": a.system_prompt,
            "model": a.model,
            "temperature": a.temperature,
            "tools_enabled": a.tools_enabled or [],
            "is_default": a.is_default,
        }
        for a in Agent.objects.all()
    ]
    tools_data = [
        {
            "slug": t.slug,
            "name": t.name,
            "description": t.description,
            "icon": t.icon,
            "is_human_in_loop": t.is_human_in_loop,
            "uses_session": t.uses_session,
            "parameters": t.parameters,
            "required": t.required,
        }
        for t in all_tools()
    ]
    app_settings = AppSettings.get_solo()
    return JsonResponse({
        "agents": agents_data,
        "models": MODEL_OPTIONS,
        "tools": tools_data,
        "playbooks": [_playbook_summary(p) for p in Playbook.objects.all()],
        "settings": {
            "max_iterations": app_settings.max_iterations,
            "massiva_workers": app_settings.massiva_workers,
        },
    })


@csrf_exempt
@require_http_methods(["POST"])
def config_settings_save(request):
    """Atualiza as configurações globais da aplicação (singleton)."""
    s = AppSettings.get_solo()
    try:
        data = json.loads(request.body)
        if "max_iterations" in data:
            try:
                n = int(data["max_iterations"])
            except (TypeError, ValueError):
                return JsonResponse(
                    {"status": "error", "message": "max_iterations inválido."},
                    status=400,
                )
            # Limites de sanidade: ao menos 1 passo; teto evita loops caros.
            s.max_iterations = max(1, min(100, n))
        if "massiva_workers" in data:
            try:
                w = int(data["massiva_workers"])
            except (TypeError, ValueError):
                return JsonResponse(
                    {"status": "error", "message": "massiva_workers inválido."},
                    status=400,
                )
            # 1–10: teto de 10 evita saturar o provider (rate limit) na
            # análise massiva. O motor também clampa por segurança.
            s.massiva_workers = max(1, min(10, w))
        s.save()
        return JsonResponse({
            "status": "success",
            "settings": {
                "max_iterations": s.max_iterations,
                "massiva_workers": s.massiva_workers,
            },
        })
    except Exception as e:
        return JsonResponse({"status": "error", "message": str(e)}, status=400)


@csrf_exempt
@require_http_methods(["POST"])
def config_agent_save(request, slug):
    """Atualiza um agente existente."""
    agent = get_object_or_404(Agent, slug=slug)
    try:
        data = json.loads(request.body)
        for field in ("name", "description", "icon", "system_prompt", "model"):
            if field in data:
                setattr(agent, field, data[field])
        if "temperature" in data:
            try:
                t = float(data["temperature"])
                agent.temperature = max(0.0, min(2.0, t))
            except (TypeError, ValueError):
                pass
        if "tools_enabled" in data and isinstance(data["tools_enabled"], list):
            agent.tools_enabled = data["tools_enabled"]
        if "is_default" in data and data["is_default"]:
            Agent.objects.exclude(id=agent.id).update(is_default=False)
            agent.is_default = True
        agent.save()
        return JsonResponse({"status": "success", "slug": agent.slug})
    except Exception as e:
        return JsonResponse({"status": "error", "message": str(e)}, status=400)


# ── Playbooks (pipelines multi-agente autorados no canvas) ───────────

def _playbook_summary(p: Playbook) -> dict:
    return {
        "id": p.id,
        "name": p.name,
        "description": p.description,
        "icon": p.icon,
        "node_count": len(p.nodes or []),
    }


def _playbook_detail(p: Playbook) -> dict:
    return {
        "id": p.id,
        "name": p.name,
        "description": p.description,
        "icon": p.icon,
        "nodes": p.nodes or [],
        "edges": p.edges or [],
        "suggestions": p.suggestions or [],
    }


def _slugify_node(name: str) -> str:
    """Slug ASCII simples a partir do nome de um nó (fallback 'no')."""
    base = re.sub(r"[^a-z0-9]+", "_", (name or "").strip().lower()).strip("_")
    return base or "no"


def _normalize_playbook_graph(payload: dict) -> tuple[dict, list]:
    """Valida e canoniza {nodes, edges, suggestions} de um playbook.

    Regras:
      - gera/mantém slugs estáveis e únicos por playbook; reescreve as arestas
        para referenciar slugs (o cliente pode mandar edges por 'id' de nó);
      - exatamente 1 nó root; root e não-root alcançáveis do root via BFS;
      - valida model ∈ MODEL_OPTIONS e tools_enabled ⊆ tools registradas.

    Retorna (graph_dict, warnings). Lança ValueError em erro fatal (→ 400).
    """
    raw_nodes = payload.get("nodes")
    if not isinstance(raw_nodes, list) or not raw_nodes:
        raise ValueError("O playbook precisa de ao menos um nó.")

    valid_tools = _valid_tool_slugs()
    warnings: list[str] = []

    # 1) Canoniza slugs. O cliente identifica cada nó por 'id' (para casar as
    #    arestas do canvas); mapeamos id → slug canônico e desduplicamos.
    id_to_slug: dict = {}
    used_slugs: set = set()
    nodes: list = []
    root_count = 0
    for i, n in enumerate(raw_nodes):
        if not isinstance(n, dict):
            raise ValueError("Nó inválido no grafo.")
        client_id = str(n.get("id") if n.get("id") is not None else n.get("slug") or i)
        slug = (n.get("slug") or "").strip() or _slugify_node(n.get("name") or "")
        # desduplica dentro do playbook
        cand, k = slug, 2
        while cand in used_slugs:
            cand = f"{slug}_{k}"
            k += 1
        slug = cand
        used_slugs.add(slug)
        id_to_slug[client_id] = slug
        id_to_slug[slug] = slug  # aresta pode já vir por slug

        model = n.get("model")
        if model not in MODEL_OPTIONS:
            raise ValueError(f"Modelo inválido no nó '{n.get('name') or slug}': {model}.")
        tools_enabled = [s for s in (n.get("tools_enabled") or []) if s in valid_tools]
        try:
            temperature = max(0.0, min(2.0, float(n.get("temperature", 0.7))))
        except (TypeError, ValueError):
            temperature = 0.7
        is_root = bool(n.get("is_root"))
        if is_root:
            root_count += 1
        canvas = n.get("canvas") if isinstance(n.get("canvas"), dict) else {}
        nodes.append({
            "slug": slug,
            "name": (str(n.get("name") or slug))[:80],
            "description": (str(n.get("description") or ""))[:240],
            "icon": (str(n.get("icon") or "🤖"))[:8] or "🤖",
            "system_prompt": str(n.get("system_prompt") or ""),
            "model": model,
            "temperature": temperature,
            "tools_enabled": tools_enabled,
            "is_root": is_root,
            "canvas": {
                "x": canvas.get("x", 0) if isinstance(canvas.get("x"), (int, float)) else 0,
                "y": canvas.get("y", 0) if isinstance(canvas.get("y"), (int, float)) else 0,
            },
        })

    if root_count != 1:
        raise ValueError(
            f"O playbook precisa de exatamente 1 nó root (encontrados {root_count})."
        )
    root_slug = next(n["slug"] for n in nodes if n["is_root"])

    # 2) Reescreve arestas para slugs; descarta as que referenciam nós inexistentes.
    node_slugs = {n["slug"] for n in nodes}
    edges: list = []
    seen_edges: set = set()
    for e in (payload.get("edges") or []):
        if not isinstance(e, dict):
            continue
        src = id_to_slug.get(str(e.get("source")))
        tgt = id_to_slug.get(str(e.get("target")))
        if src not in node_slugs or tgt not in node_slugs or src == tgt:
            continue
        key = (src, tgt)
        if key in seen_edges:
            continue
        seen_edges.add(key)
        edges.append({"source": src, "target": tgt})

    # 3) Alcançabilidade: todo nó não-root precisa ser alcançável do root.
    adjacency: dict = {}
    for e in edges:
        adjacency.setdefault(e["source"], []).append(e["target"])
    reachable = {root_slug}
    frontier = [root_slug]
    while frontier:
        cur = frontier.pop()
        for nxt in adjacency.get(cur, []):
            if nxt not in reachable:
                reachable.add(nxt)
                frontier.append(nxt)
    unreachable = node_slugs - reachable
    if unreachable:
        nomes = ", ".join(sorted(unreachable))
        raise ValueError(
            f"Nós inalcançáveis a partir do root: {nomes}. Conecte-os com arestas."
        )

    # 4) Avisos não-fatais: call_agent habilitado sem arestas de saída.
    for n in nodes:
        if "call_agent" in n["tools_enabled"] and not adjacency.get(n["slug"]):
            warnings.append(
                f"Nó '{n['name']}' tem call_agent habilitado mas nenhuma aresta de saída."
            )

    suggestions = []
    for s in (payload.get("suggestions") or []):
        if not isinstance(s, dict):
            continue
        title = (str(s.get("title") or "")).strip()[:80]
        text = (str(s.get("text") or "")).strip()[:500]
        if text:
            suggestions.append({"title": title, "text": text})

    return {"nodes": nodes, "edges": edges, "suggestions": suggestions}, warnings


def _apply_playbook_fields(p: Playbook, data: dict) -> list:
    """Valida e aplica os campos de um Playbook (sem salvar). Retorna warnings."""
    name = (str(data.get("name") or "")).strip()
    if not name:
        raise ValueError("Nome é obrigatório.")
    p.name = name[:120]
    p.description = (str(data.get("description") or "")).strip()[:240]
    p.icon = (str(data.get("icon") or "📘")).strip()[:8] or "📘"
    graph, warnings = _normalize_playbook_graph(data)
    p.nodes = graph["nodes"]
    p.edges = graph["edges"]
    p.suggestions = graph["suggestions"]
    return warnings


@require_GET
def playbook_list(request):
    """Lista todos os playbooks cadastrados (resumo p/ pickers)."""
    items = [_playbook_summary(p) for p in Playbook.objects.all()]
    return JsonResponse({"status": "success", "playbooks": items})


@require_GET
def playbook_detail(request, pb_id):
    """Retorna o grafo completo de um playbook (para o editor de canvas)."""
    p = get_object_or_404(Playbook, id=pb_id)
    return JsonResponse({"status": "success", "playbook": _playbook_detail(p)})


@csrf_exempt
@require_http_methods(["POST"])
def playbook_create(request):
    """Cria um novo playbook."""
    try:
        data = json.loads(request.body or "{}")
        p = Playbook()
        warnings = _apply_playbook_fields(p, data)
        p.save()
        return JsonResponse({
            "status": "success",
            "playbook": _playbook_detail(p),
            "warnings": warnings,
        })
    except ValueError as e:
        return JsonResponse({"status": "error", "message": str(e)}, status=400)
    except Exception as e:
        return JsonResponse({"status": "error", "message": str(e)}, status=400)


@csrf_exempt
@require_http_methods(["POST"])
def playbook_update(request, pb_id):
    """Atualiza um playbook existente."""
    p = get_object_or_404(Playbook, id=pb_id)
    try:
        data = json.loads(request.body or "{}")
        warnings = _apply_playbook_fields(p, data)
        p.save()
        return JsonResponse({
            "status": "success",
            "playbook": _playbook_detail(p),
            "warnings": warnings,
        })
    except ValueError as e:
        return JsonResponse({"status": "error", "message": str(e)}, status=400)
    except Exception as e:
        return JsonResponse({"status": "error", "message": str(e)}, status=400)


@csrf_exempt
@require_http_methods(["POST", "DELETE"])
def playbook_delete(request, pb_id):
    """Exclui um playbook. Conversas vinculadas revertem para o agente global
    no próximo turno (FK SET_NULL)."""
    p = get_object_or_404(Playbook, id=pb_id)
    try:
        p.delete()
        return JsonResponse({"status": "success"})
    except Exception as e:
        return JsonResponse({"status": "error", "message": str(e)}, status=400)


@csrf_exempt
@require_http_methods(["POST"])
def conversation_playbook_save(request, conv_id):
    """Vincula (ou desvincula, com playbook_id=null) um playbook à conversa."""
    conv = get_object_or_404(Conversation, id=conv_id)
    try:
        data = json.loads(request.body or "{}")
        pb_id = data.get("playbook_id")
        if pb_id is None:
            conv.playbook = None
        else:
            conv.playbook = get_object_or_404(Playbook, id=pb_id)
        conv.save(update_fields=["playbook", "updated_at"])
        return JsonResponse({"status": "success", "conversation": _conversation_summary(conv)})
    except Exception as e:
        return JsonResponse({"status": "error", "message": str(e)}, status=400)


# Sentinela: distingue "playbook_id não enviado" (não mexe no vínculo) de
# "playbook_id: null" (desvincula explicitamente).
_UNSET = object()


def _apply_playbook(conv, playbook_id) -> None:
    """Aplica a vinculação de playbook vinda do payload do chat.

    playbook_id == _UNSET → não mexe (o cliente não mandou o campo).
    playbook_id is None    → desvincula.
    playbook_id: int       → vincula (ignora silenciosamente id inexistente).
    """
    if playbook_id is _UNSET:
        return
    if playbook_id is None:
        if conv.playbook_id is not None:
            conv.playbook = None
            conv.save(update_fields=["playbook", "updated_at"])
        return
    pb = Playbook.objects.filter(id=playbook_id).first()
    if pb is not None and conv.playbook_id != pb.id:
        conv.playbook = pb
        conv.save(update_fields=["playbook", "updated_at"])


def _playbook_runtime(conv):
    """Se a conversa tem playbook, retorna (root_RuntimeAgent, session_extras).

    Senão (None, {}). As extras são chaves '__' transientes (o _persist_turn as
    remove) que o call_agent/run_agent leem para resolver a delegação isolada.
    """
    pb = conv.playbook
    if pb is None or not pb.nodes:
        return None, {}
    from .ai_service import _node_to_runtime_agent
    nodes = {n["slug"]: n for n in pb.nodes if isinstance(n, dict) and n.get("slug")}
    root_slug = next((n["slug"] for n in pb.nodes if n.get("is_root")), None)
    if root_slug is None or root_slug not in nodes:
        return None, {}
    adjacency: dict = {}
    for e in (pb.edges or []):
        if isinstance(e, dict) and e.get("source") and e.get("target"):
            adjacency.setdefault(e["source"], []).append(e["target"])
    extras = {
        "__playbook_nodes": nodes,
        "__playbook_adjacency": adjacency,
        "__playbook_root": root_slug,
        "__current_node": root_slug,
    }
    return _node_to_runtime_agent(nodes[root_slug]), extras


# ── Agente da sessão (criado só para uma conversa) ───────────────────

_VALID_TOOL_SLUGS = None  # cache preguiçoso


def _valid_tool_slugs() -> set:
    global _VALID_TOOL_SLUGS
    if _VALID_TOOL_SLUGS is None:
        _VALID_TOOL_SLUGS = {t.slug for t in all_tools()}
    return _VALID_TOOL_SLUGS


def _apply_session_agent_fields(sa: SessionAgent, data: dict) -> None:
    """Valida e aplica os campos do formulário num SessionAgent (sem salvar)."""
    if "name" in data:
        sa.name = (str(data["name"]).strip() or "Meu agente")[:80]
    if "icon" in data and data["icon"]:
        sa.icon = str(data["icon"])[:8]
    if "system_prompt" in data:
        sa.system_prompt = str(data["system_prompt"] or "")
    if "guardrails" in data:
        sa.guardrails = str(data["guardrails"] or "")
    if "model" in data and data["model"] in MODEL_OPTIONS:
        sa.model = data["model"]
    if "temperature" in data:
        try:
            sa.temperature = max(0.0, min(2.0, float(data["temperature"])))
        except (TypeError, ValueError):
            pass
    if "tools_enabled" in data and isinstance(data["tools_enabled"], list):
        valid = _valid_tool_slugs()
        sa.tools_enabled = [s for s in data["tools_enabled"] if s in valid]
    if "documents" in data and isinstance(data["documents"], list):
        sa.documents = _merge_documents(sa.documents or [], data["documents"])


def _merge_documents(existing: list, incoming: list) -> list:
    """Reconcilia a lista de documentos vinda do formulário com a do banco.

    O frontend só recebe o RESUMO dos docs já salvos (sem o markdown pesado).
    No save ele devolve a lista atual: docs antigos vêm como referência (só
    `filename`, sem markdown) e docs novos vêm completos (com markdown, do
    endpoint de upload). Aqui:
      - doc com `markdown` → usa como veio (recém-anexado);
      - doc sem `markdown` → recupera o conteúdo do banco pelo filename;
      - doc removido na UI → simplesmente não aparece em `incoming` → cai fora.
    """
    by_name = {d.get("filename"): d for d in existing if isinstance(d, dict)}
    merged = []
    for item in incoming:
        if not isinstance(item, dict):
            continue
        filename = item.get("filename")
        if item.get("markdown"):
            md = str(item["markdown"])
            merged.append({
                "filename": filename or "documento",
                "markdown": md,
                "char_count": len(md),
                "page_count": item.get("page_count"),
            })
        elif filename in by_name:
            merged.append(by_name[filename])
    return merged


@require_GET
def session_agent_detail(request, conv_id):
    """Retorna o agente da sessão de uma conversa (ou 404 se não existir)."""
    conv = get_object_or_404(Conversation, id=conv_id)
    sa = getattr(conv, "session_agent", None)
    if sa is None:
        return JsonResponse({"status": "error", "message": "Sem agente de sessão."}, status=404)
    return JsonResponse({"status": "success", "agent": _session_agent_payload(sa)})


@csrf_exempt
@require_http_methods(["POST"])
def session_agent_save(request, conv_id):
    """Cria ou atualiza o agente da sessão de uma conversa (upsert)."""
    conv = get_object_or_404(Conversation, id=conv_id)
    try:
        data = json.loads(request.body or "{}")
        sa, _created = SessionAgent.objects.get_or_create(conversation=conv)
        _apply_session_agent_fields(sa, data)
        sa.save()
        return JsonResponse({"status": "success", "agent": _session_agent_payload(sa)})
    except Exception as e:
        return JsonResponse({"status": "error", "message": str(e)}, status=400)


@csrf_exempt
@require_http_methods(["DELETE"])
def session_agent_delete(request, conv_id):
    """Remove o agente da sessão de uma conversa."""
    conv = get_object_or_404(Conversation, id=conv_id)
    SessionAgent.objects.filter(conversation=conv).delete()
    return JsonResponse({"status": "success"})


@csrf_exempt
@require_http_methods(["POST"])
def session_agent_create_conversation(request):
    """Cria uma conversa nova já com um agente de sessão.

    Usado quando o usuário monta o agente ANTES de mandar a 1ª mensagem
    (chat novo). A conversa nasce associada ao agente global default (o
    orquestrador), que então poderá delegar para o agente da sessão.
    """
    try:
        data = json.loads(request.body or "{}")
        agent = Agent.objects.filter(is_default=True).first() or Agent.objects.first()
        if agent is None:
            return JsonResponse(
                {"status": "error", "message": "Nenhum agente cadastrado. Acesse Configurações."},
                status=400,
            )
        nome = (str(data.get("name") or "").strip() or "Meu agente")
        conv = Conversation.objects.create(title=f"🤖 {nome}"[:120], agent=agent)
        sa = SessionAgent(conversation=conv)
        _apply_session_agent_fields(sa, data)
        sa.save()
        return JsonResponse({
            "status": "success",
            "conversation_id": conv.id,
            "agent": _session_agent_payload(sa),
        })
    except Exception as e:
        return JsonResponse({"status": "error", "message": str(e)}, status=400)


# Extensões aceitas como documento do agente da sessão.
SA_DOC_EXTS = {".pdf", ".txt", ".md", ".docx", ".doc"}


@csrf_exempt
@require_http_methods(["POST"])
def session_agent_extract_document(request):
    """Extrai o texto (markdown) de um documento para anexar ao agente.

    NÃO persiste nada — apenas devolve o conteúdo extraído. O frontend guarda
    o doc na lista do formulário e o conteúdo só é gravado quando o usuário
    salva o agente (em session_agent_save). Assim funciona tanto num chat já
    existente quanto num chat novo (ainda sem conversa criada).
    """
    f = request.FILES.get("file")
    if not f:
        return JsonResponse({"status": "error", "message": "Nenhum arquivo enviado."}, status=400)

    ext = os.path.splitext(f.name)[1].lower()
    if ext not in SA_DOC_EXTS:
        aceitas = ", ".join(sorted(SA_DOC_EXTS))
        return JsonResponse(
            {"status": "error", "message": f"Formato não suportado. Aceitos: {aceitas}."},
            status=400,
        )

    try:
        if ext in (".txt", ".md"):
            # Texto puro: lê direto, sem docling (decodifica com fallback).
            f.seek(0)
            raw = f.read()
            markdown = raw.decode("utf-8", errors="replace") if isinstance(raw, bytes) else str(raw)
            page_count = None
        else:
            markdown, page_count = _docling_to_markdown(f, f.name)
    except RuntimeError as e:
        return JsonResponse({"status": "error", "message": str(e)}, status=500)
    except Exception as e:
        return JsonResponse(
            {"status": "error", "message": f"Falha ao extrair documento ({type(e).__name__}): {e}"},
            status=400,
        )

    markdown = markdown or ""
    if not markdown.strip():
        return JsonResponse(
            {"status": "error", "message": "Não foi possível extrair texto do documento."},
            status=400,
        )

    return JsonResponse({
        "status": "success",
        "document": {
            "filename": f.name,
            "markdown": markdown,
            "char_count": len(markdown),
            "page_count": page_count,
        },
    })


# ── Upload de tabelas (CSV/XLSX) ─────────────────────────────────────

def _read_table(file_obj, filename: str):
    """Lê CSV, XLSX ou Parquet a partir de um InMemoryUploadedFile, retorna DataFrame."""
    import pandas as pd

    ext = os.path.splitext(filename)[1].lower()
    if ext == ".csv":
        try:
            return pd.read_csv(file_obj)
        except UnicodeDecodeError:
            file_obj.seek(0)
            return pd.read_csv(file_obj, encoding="latin-1")
        except Exception:
            file_obj.seek(0)
            return pd.read_csv(file_obj, sep=";")
    if ext in (".xlsx", ".xls"):
        if hasattr(file_obj, "seek"):
            file_obj.seek(0)
        xls = pd.ExcelFile(file_obj)
        if len(xls.sheet_names) <= 1:
            df = pd.read_excel(xls, sheet_name=0)
            df.attrs["sheet_name"] = xls.sheet_names[0] if xls.sheet_names else None
            df.attrs["available_sheets"] = list(xls.sheet_names)
            df.attrs["sheet_datasets"] = {
                xls.sheet_names[0]: df.copy()
            } if xls.sheet_names else {}
            return df

        def sheet_score(df):
            non_empty_rows = df.dropna(how="all").shape[0]
            non_empty_cols = df.dropna(axis=1, how="all").shape[1]
            non_null_cells = int(df.notna().sum().sum())
            return (non_empty_rows, non_empty_cols, non_null_cells)

        filename_base = os.path.splitext(filename)[0].lower()
        preferred_keywords = [filename_base, "accounts", "account", "dados", "data", "sheet", "planilha"]

        best_sheet = None
        best_sheet_name = None
        best_score = (-1, -1, -1)
        best_preferred = None
        # Mantém cada aba legível na sessão. Antes o upload escolhia uma aba
        # "principal" e perdia as demais, impedindo o agente de abrir uma
        # sheet como Transactions depois do upload.
        sheet_datasets = {}
        for sheet_name in xls.sheet_names:
            try:
                df = xls.parse(sheet_name=sheet_name)
            except Exception:
                continue
            if df is None:
                continue
            sheet_datasets[str(sheet_name)] = df
            score = sheet_score(df)
            if score[0] == 0 or score[1] == 0:
                continue

            normalized_sheet_name = str(sheet_name).strip().lower()
            is_preferred = any(keyword in normalized_sheet_name for keyword in preferred_keywords if keyword)
            if is_preferred:
                if best_preferred is None or score > best_preferred[0]:
                    best_preferred = (score, sheet_name, df)
            if score > best_score:
                best_score = score
                best_sheet_name = sheet_name
                best_sheet = df

        if best_preferred is not None:
            score, sheet_name, df = best_preferred
            df.attrs["sheet_name"] = sheet_name
            df.attrs["available_sheets"] = list(xls.sheet_names)
            # Cópias evitam que o DataFrame ativo referencie a si próprio
            # dentro de ``attrs`` (o que causa RecursionError no to_json).
            df.attrs["sheet_datasets"] = {
                name: frame.copy(deep=True) for name, frame in sheet_datasets.items()
            }
            return df
        if best_sheet is not None:
            best_sheet.attrs["sheet_name"] = best_sheet_name
            best_sheet.attrs["available_sheets"] = list(xls.sheet_names)
            best_sheet.attrs["sheet_datasets"] = {
                name: frame.copy(deep=True) for name, frame in sheet_datasets.items()
            }
            return best_sheet
        df = pd.read_excel(xls, sheet_name=0)
        df.attrs["sheet_name"] = xls.sheet_names[0]
        df.attrs["available_sheets"] = list(xls.sheet_names)
        df.attrs["sheet_datasets"] = {xls.sheet_names[0]: df.copy()}
        return df
    if ext == ".parquet":
        return pd.read_parquet(file_obj)
    raise ValueError(f"Extensão não suportada: {ext}. Envie .csv, .xlsx, .xls ou .parquet.")


def _df_to_records_jsonsafe(df) -> list[dict]:
    """
    Converte um DataFrame em lista de dicts 100% JSON-safe.

    df.to_dict() preserva NaN/NaT/numpy types — o json.dumps do Python emite
    `NaN` literal, que não é JSON válido e o SQLite rejeita via CHECK constraint.
    Roteamos pelo to_json do pandas (que mapeia NaN→null, datas→ISO).
    """
    import json as _json
    raw = df.to_json(orient="records", date_format="iso", default_handler=str)
    return _json.loads(raw)


def _llm_summary_for_table(filename: str, df) -> str:
    """Texto curto que vai pro modelo: filename + shape + colunas + dtypes + 3 linhas."""
    import json as _json

    sample = _df_to_records_jsonsafe(df.head(PREVIEW_ROWS_FOR_LLM))
    dtypes = {c: str(df[c].dtype) for c in df.columns}
    payload = {
        "arquivo_carregado": filename,
        "linhas": int(len(df)),
        "colunas": list(df.columns),
        "dtypes": dtypes,
        "amostra_3_linhas": sample,
    }
    if df.attrs.get("sheet_name"):
        payload["sheet_name"] = df.attrs["sheet_name"]
    available_sheets = df.attrs.get("available_sheets") or []
    if available_sheets:
        payload["sheets_disponiveis"] = available_sheets

    sheet_line = ""
    if df.attrs.get("sheet_name"):
        sheet_line = f" O Excel foi carregado da planilha **{df.attrs['sheet_name']}**."
    other_sheets_line = ""
    if available_sheets:
        dataset_names = ", ".join(
            f"`{filename}::{sheet}`" for sheet in available_sheets
        )
        other_sheets_line = (
            f" As abas disponíveis neste Excel são: {', '.join(map(str, available_sheets))}."
            " Todas foram mantidas como datasets nomeados na sessão "
            f"({dataset_names}); use `executar_pandas` com "
            "`result_df = load_dataset('<arquivo>::<aba>')` para trocar a aba corrente."
        )

    return (
        f"Anexei o arquivo **{filename}** ({len(df)} linhas × {len(df.columns)} colunas)."
        f"{sheet_line}{other_sheets_line} "
        f"O dataset completo está disponível na sessão para as tools de análise — "
        f"você só vê os metadados e 3 linhas abaixo:\n\n```json\n"
        + _json.dumps(payload, ensure_ascii=False, indent=2, default=str)
        + "\n```"
    )


def _bundled_rapidocr_models() -> dict | None:
    """
    Caminhos dos modelos ONNX que vêm embutidos no pacote
    rapidocr_onnxruntime (det/rec/cls). Apontar explicitamente para eles
    garante OCR 100% offline — sem o download de modelscope.cn que o
    ambiente bloqueia. Retorna None se o pacote não estiver instalado.
    """
    try:
        import rapidocr_onnxruntime as _ror
    except ImportError:
        return None

    models_dir = os.path.join(os.path.dirname(_ror.__file__), "models")
    paths = {
        "det_model_path": os.path.join(models_dir, "ch_PP-OCRv4_det_infer.onnx"),
        "rec_model_path": os.path.join(models_dir, "ch_PP-OCRv4_rec_infer.onnx"),
        "cls_model_path": os.path.join(models_dir, "ch_ppocr_mobile_v2.0_cls_infer.onnx"),
    }
    if all(os.path.exists(p) for p in paths.values()):
        return paths
    return None


def _build_docling_converter():
    """
    Monta o DocumentConverter com OCR habilitado via backend onnxruntime e
    modelos locais (sem download). Aplica as opções a PDF e imagens.

    Quando `arquivos_suporte/docling/` contém os modelos de layout/tabela
    baixados manualmente, passa `artifacts_path` para o pipeline — assim
    o docling não tenta nenhum request à HuggingFace.
    """
    from django.conf import settings
    from docling.datamodel.base_models import InputFormat
    from docling.datamodel.pipeline_options import PdfPipelineOptions, RapidOcrOptions
    from docling.document_converter import (
        DocumentConverter,
        PdfFormatOption,
        ImageFormatOption,
    )

    from auditor.proxy_config import docling_artifacts_path

    bundled = _bundled_rapidocr_models()
    ocr_options = RapidOcrOptions(backend="onnxruntime", **(bundled or {}))

    artifacts = docling_artifacts_path(settings.BASE_DIR)
    pipeline_kwargs = {"do_ocr": True, "ocr_options": ocr_options}
    if artifacts is not None:
        pipeline_kwargs["artifacts_path"] = str(artifacts)
    pipeline_options = PdfPipelineOptions(**pipeline_kwargs)

    return DocumentConverter(
        format_options={
            InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options),
            InputFormat.IMAGE: ImageFormatOption(pipeline_options=pipeline_options),
        }
    )


def _docling_to_markdown(file_obj, filename: str) -> tuple[str, int | None]:
    """
    Converte PDF/DOCX/imagem em markdown via docling (com OCR para imagens
    e PDFs escaneados). Retorna (markdown, page_count).
    """
    import tempfile

    try:
        from docling.document_converter import DocumentConverter
    except ImportError as e:
        raise RuntimeError(
            "docling não está instalado. Rode:\n"
            "pip install --index-url https://artifactory.prod.aws.cloud.ihf/artifactory/api/pypi/python-devel/simple "
            "--trusted-host artifactory.prod.aws.cloud.ihf docling rapidocr_onnxruntime"
        ) from e
    except ModuleNotFoundError as e:
        # Transformers não instalado ou corrompido - modo fallback
        if "transformers" in str(e):
            raise RuntimeError(
                "transformers não pôde ser instalado devido ao limite de caminho longo do Windows.\n\n"
                "Abra PowerShell como Administrador e execute:\n"
                "New-ItemProperty -Path 'HKLM:\\SYSTEM\\CurrentControlSet\\Control\\FileSystem' "
                "-Name 'LongPathsEnabled' -Value 1 -PropertyType DWORD -Force\n\n"
                "Depois reinicie o computador e execute:\n"
                "pip install --index-url https://artifactory.prod.aws.cloud.ihf/artifactory/api/pypi/python-devel/simple "
                "--trusted-host artifactory.prod.aws.cloud.ihf transformers==4.46.3 --no-deps"
            ) from e
        raise RuntimeError(f"Erro ao importar docling: {e}") from e
    except Exception as e:
        raise RuntimeError(
            f"Erro ao importar docling ({type(e).__name__}): {e}"
        ) from e

    ext = os.path.splitext(filename)[1].lower()
    file_obj.seek(0)
    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        for chunk in file_obj.chunks():
            tmp.write(chunk)
        tmp_path = tmp.name

    try:
        try:
            converter = _build_docling_converter()
        except Exception as e:
            # Fallback: converter padrão (OCR pode ficar indisponível, mas
            # PDFs/DOCX com texto nativo continuam funcionando).
            try:
                converter = DocumentConverter()
            except Exception as fallback_error:
                raise RuntimeError(
                    f"Falha ao criar DocumentConverter ({type(fallback_error).__name__}): {fallback_error}. "
                    "Verifique a instalação do torch: pip list | grep -E '(torch|transformers)'"
                ) from fallback_error
        result = converter.convert(tmp_path)
        md = result.document.export_to_markdown()
        page_count = None
        pages = getattr(result.document, "pages", None)
        if pages is not None:
            try:
                page_count = len(pages)
            except TypeError:
                page_count = None
        return md, page_count
    except RuntimeError:
        raise
    except Exception as e:
        raise RuntimeError(
            f"Falha ao extrair documento ({type(e).__name__}): {e}"
        ) from e
    finally:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass


def _pptx_to_markdown(file_obj, filename: str) -> tuple[str, int | None]:
    """
    Extrai o conteúdo de um .pptx via python-pptx com mais fidelidade que o
    docling: percorre slide a slide capturando título, textos das formas,
    tabelas (em markdown) e as notas do apresentador (speaker notes).

    Retorna (markdown, page_count) onde page_count = nº de slides. Mantém a
    mesma forma de retorno de `_docling_to_markdown` para reaproveitar o fluxo
    de upload e as tools de documento existentes.
    """
    try:
        from pptx import Presentation
        from pptx.util import Emu  # noqa: F401 (garante que o pacote está OK)
    except ImportError as e:
        raise RuntimeError(
            "python-pptx não está instalado. Rode: pip install python-pptx"
        ) from e

    file_obj.seek(0)
    try:
        prs = Presentation(file_obj)
    except Exception as e:
        raise RuntimeError(
            f"Falha ao abrir o PPTX ({type(e).__name__}): {e}"
        ) from e

    def _table_to_md(table) -> str:
        rows = []
        for row in table.rows:
            cells = [(c.text or "").strip().replace("\n", " ") for c in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
        if not rows:
            return ""
        # cabeçalho = primeira linha; separador logo abaixo
        n_cols = rows[0].count("|") - 1
        sep = "| " + " | ".join(["---"] * max(n_cols, 1)) + " |"
        return "\n".join([rows[0], sep, *rows[1:]])

    partes: list[str] = []
    slides = list(prs.slides)
    for idx, slide in enumerate(slides, start=1):
        # Título do slide (placeholder de título, quando existir)
        titulo = ""
        try:
            if slide.shapes.title is not None and slide.shapes.title.text:
                titulo = slide.shapes.title.text.strip()
        except (AttributeError, ValueError):
            titulo = ""

        cabecalho = f"## Slide {idx}" + (f" — {titulo}" if titulo else "")
        blocos: list[str] = [cabecalho]

        for shape in slide.shapes:
            # Não repetir o título já usado no cabeçalho
            if titulo and getattr(shape, "text", None) and shape.text.strip() == titulo:
                continue
            if getattr(shape, "has_table", False) and shape.has_table:
                md_tab = _table_to_md(shape.table)
                if md_tab:
                    blocos.append(md_tab)
            elif getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                texto = "\n".join(
                    p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()
                )
                if texto:
                    blocos.append(texto)

        # Notas do apresentador
        try:
            if slide.has_notes_slide:
                notas = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notas:
                    blocos.append(f"**Notas do apresentador:**\n{notas}")
        except (AttributeError, ValueError):
            pass

        partes.append("\n\n".join(blocos))

    markdown = "\n\n".join(partes).strip()
    return markdown, len(slides)


def _llm_summary_for_document(filename: str, markdown: str, page_count: int | None) -> str:
    """Texto curto que vai pro modelo: nome, tamanho e início do conteúdo."""
    preview = markdown[:1500]
    if len(markdown) > 1500:
        preview += "\n…"
    pages_part = f", {page_count} página(s)" if page_count else ""
    return (
        f"Anexei o documento **{filename}** "
        f"({len(markdown):,} caracteres{pages_part}). "
        f"O conteúdo extraído (markdown) está disponível na sessão para as "
        f"tools de análise de documento. Início do conteúdo:\n\n"
        f"```markdown\n{preview}\n```"
    )


def _handle_document_upload(f, *, conv_id, agent_slug, user_note: str):
    """Fluxo de upload para PDF/DOCX/imagem via docling.

    Exceção: .pptx usa python-pptx (extração por slide, com notas e tabelas),
    que preserva mais estrutura que o markdown achatado do docling. O formato
    antigo .ppt (binário) o python-pptx não abre, então continua no docling.
    """
    ext = os.path.splitext(f.name)[1].lower()
    try:
        if ext == ".pptx":
            markdown, page_count = _pptx_to_markdown(f, f.name)
        else:
            markdown, page_count = _docling_to_markdown(f, f.name)
    except RuntimeError as e:
        import traceback
        traceback.print_exc()
        return JsonResponse({"status": "error", "message": str(e)}, status=500)
    except Exception as e:
        import traceback
        traceback.print_exc()
        return JsonResponse(
            {
                "status": "error",
                "message": f"Falha ao extrair documento ({type(e).__name__}): {e}",
            },
            status=400,
        )

    # Recupera/cria a conversa
    conv: Conversation | None = None
    if conv_id:
        conv = Conversation.objects.filter(id=conv_id).first()
    if conv is None:
        agent = None
        if agent_slug:
            agent = Agent.objects.filter(slug=agent_slug).first()
        if agent is None:
            agent = (
                Agent.objects.filter(slug="analista_documentos").first()
                or Agent.objects.filter(is_default=True).first()
                or Agent.objects.first()
            )
        if agent is None:
            return JsonResponse(
                {"status": "error", "message": "Nenhum agente cadastrado."},
                status=400,
            )
        title = f"📄 {f.name}"[:120]
        conv = Conversation.objects.create(title=title, agent=agent)

    # Salva o documento na sessão
    state = dict(conv.state or {})
    state["documento_atual"] = {
        "filename": f.name,
        "markdown": markdown,
        "char_count": len(markdown),
        "page_count": page_count,
    }
    conv.state = state
    conv.save()

    summary = _llm_summary_for_document(f.name, markdown, page_count)
    if user_note:
        summary = f"{user_note}\n\n{summary}"

    attachment = {
        "kind": "document",
        "filename": f.name,
        "char_count": len(markdown),
        "page_count": page_count,
        "preview": markdown[:4000],
    }

    msg = Message.objects.create(
        conversation=conv,
        role="user",
        content=summary,
        attachment=attachment,
    )

    return JsonResponse({
        "status": "success",
        "conversation_id": conv.id,
        "conversation_title": conv.title,
        "agent_slug": conv.agent.slug if conv.agent else None,
        "message": _message_payload(msg),
    })


def _extract_pdf_text(file_obj) -> str:
    """Extrai texto de um PDF via PyMuPDF (sem OCR)."""
    try:
        import fitz  # PyMuPDF
        data = file_obj.read() if hasattr(file_obj, "read") else file_obj
        doc = fitz.open(stream=data, filetype="pdf")
        pages = [page.get_text() for page in doc]
        doc.close()
        return "\n\n".join(pages).strip()
    except Exception as e:
        return f"[erro ao extrair: {e}]"


def _extract_txt_text(file_obj) -> str:
    data = file_obj.read() if hasattr(file_obj, "read") else file_obj
    if isinstance(data, bytes):
        return data.decode("utf-8", errors="replace").strip()
    return data.strip()


def _persist_batch_docs(registros, erros, conv_id, agent_slug, user_note):
    """Monta o dataset [nome_arquivo, conteudo_extraido], salva na sessão da
    conversa e cria a mensagem de usuário com o card-tabela. Retorna o payload
    final (mesmo formato do antigo JsonResponse de sucesso)."""
    import pandas as pd
    df = pd.DataFrame(registros, columns=["nome_arquivo", "conteudo_extraido"])

    conv: Conversation | None = None
    if conv_id:
        conv = Conversation.objects.filter(id=conv_id).first()
    if conv is None:
        agent = None
        if agent_slug:
            agent = Agent.objects.filter(slug=agent_slug).first()
        if agent is None:
            agent = Agent.objects.filter(is_default=True).first() or Agent.objects.first()
        if agent is None:
            raise RuntimeError("Nenhum agente cadastrado.")
        title = f"📁 {len(registros)} documentos"[:120]
        conv = Conversation.objects.create(title=title, agent=agent)

    state = dict(conv.state or {})
    state["athena_last_result"]  = _df_to_records_jsonsafe(df)
    state["athena_last_columns"] = list(df.columns)
    state["athena_last_source"]  = {
        "kind": "batch_docs",
        "count": len(registros),
        "filenames": [r["nome_arquivo"] for r in registros],
    }
    conv.state = state
    conv.save()

    aviso_erros = ""
    if erros:
        aviso_erros = f"\n\n⚠️ {len(erros)} arquivo(s) ignorado(s) (extensão não suportada): {', '.join(erros)}"

    summary = (
        f"{user_note + chr(10) + chr(10) if user_note else ''}"
        f"📁 **{len(registros)} documento(s) carregado(s)** e extraídos para dataset.\n"
        f"Colunas: `nome_arquivo`, `conteudo_extraido`.\n"
        f"Arquivos: {', '.join(r['nome_arquivo'] for r in registros[:10])}"
        f"{'…' if len(registros) > 10 else ''}"
        f"{aviso_erros}"
    )

    preview = _df_to_records_jsonsafe(df.head(PREVIEW_ROWS_FOR_UI))
    attachment = {
        "kind": "table",
        "filename": f"{len(registros)} documentos",
        "rows": len(df),
        "columns": list(df.columns),
        "dtypes": {c: str(df[c].dtype) for c in df.columns},
        "preview": preview,
        "preview_rows": len(preview),
        "truncated": False,
    }

    msg = Message.objects.create(
        conversation=conv,
        role="user",
        content=summary,
        attachment=attachment,
    )

    return {
        "status": "success",
        "conversation_id": conv.id,
        "conversation_title": conv.title,
        "agent_slug": conv.agent.slug if conv.agent else None,
        "message": _message_payload(msg),
        "count": len(registros),
        "filenames": [r["nome_arquivo"] for r in registros],
    }


@csrf_exempt
@require_http_methods(["POST"])
def upload_batch_docs(request):
    """
    Recebe múltiplos PDFs/TXTs, extrai o texto de cada um com PyMuPDF e salva
    na sessão como dataset [nome_arquivo, conteudo_extraido].

    Responde via SSE (text/event-stream): emite um evento 'progress' por
    arquivo extraído ("extraindo N de TOTAL…") e um evento final 'done' com o
    mesmo payload do upload de tabela. O frontend consome com consumeStream().
    """
    files = request.FILES.getlist("files")
    if not files:
        return JsonResponse({"status": "error", "message": "Nenhum arquivo enviado."}, status=400)

    if len(files) > BATCH_MAX_FILES:
        return JsonResponse(
            {"status": "error",
             "message": f"Máximo de {BATCH_MAX_FILES} arquivos por vez (você enviou {len(files)})."},
            status=400,
        )

    conv_id    = request.POST.get("conversation_id")
    agent_slug = request.POST.get("agent_slug")
    user_note  = (request.POST.get("note") or "").strip()

    # Lê o conteúdo em memória AGORA — os arquivos temporários da request são
    # fechados ao fim da view, mas o generator do StreamingHttpResponse roda
    # depois. Guardamos (nome, ext, bytes) para extrair dentro do stream.
    arquivos = []
    for f in files:
        arquivos.append((f.name, os.path.splitext(f.name)[1].lower(), f.read()))

    total = len(arquivos)

    def stream():
        yield ": stream start\n\n"
        registros = []
        erros = []
        try:
            for i, (nome, ext, data) in enumerate(arquivos, start=1):
                yield "data: " + json.dumps(
                    {"type": "progress", "done": i, "total": total, "filename": nome},
                    ensure_ascii=False,
                ) + "\n\n"

                if ext not in BATCH_PDF_EXTS:
                    erros.append(nome)
                    continue
                if ext == ".pdf":
                    conteudo = _extract_pdf_text(data)
                else:
                    conteudo = _extract_txt_text(data)
                registros.append({"nome_arquivo": nome, "conteudo_extraido": conteudo})

            if not registros:
                exts = sorted(BATCH_PDF_EXTS)
                yield "data: " + json.dumps(
                    {"type": "error", "message": f"Nenhum arquivo suportado. Use: {exts}"},
                    ensure_ascii=False,
                ) + "\n\n"
                return

            payload = _persist_batch_docs(registros, erros, conv_id, agent_slug, user_note)
            yield "data: " + json.dumps({"type": "done", "payload": payload}, ensure_ascii=False) + "\n\n"
        except Exception as e:
            import traceback
            traceback.print_exc()
            yield "data: " + json.dumps(
                {"type": "error", "message": f"{type(e).__name__}: {e}"},
                ensure_ascii=False,
            ) + "\n\n"
        finally:
            from django.db import connection
            connection.close()

    response = StreamingHttpResponse(stream(), content_type="text/event-stream")
    response["Cache-Control"] = "no-cache"
    response["X-Accel-Buffering"] = "no"
    return response


@csrf_exempt
@require_http_methods(["POST"])
def upload_table(request):
    """
    Recebe um arquivo (CSV/XLSX → tabela; PDF/DOCX/PPT/imagem → documento),
    carrega para a sessão da conversa e cria uma mensagem de usuário com o resumo
    (apenas metadados+amostra vão para o LLM).

    Todo o corpo está envolvido em try/except para garantir resposta JSON
    mesmo em erros inesperados — caso contrário o Django devolveria HTML
    (página de debug) e o frontend quebra ao tentar parsear como JSON.
    """
    try:
        f = request.FILES.get("file")
        if not f:
            return JsonResponse({"status": "error", "message": "Arquivo não enviado."}, status=400)
        if f.size > UPLOAD_MAX_BYTES:
            mb = UPLOAD_MAX_BYTES // (1024 * 1024)
            return JsonResponse(
                {"status": "error", "message": f"Arquivo maior que {mb}MB."},
                status=400,
            )

        conv_id = request.POST.get("conversation_id")
        agent_slug = request.POST.get("agent_slug")
        user_note = (request.POST.get("note") or "").strip()

        ext = os.path.splitext(f.name)[1].lower()
        is_table = ext in TABLE_EXTS
        is_document = ext in DOCUMENT_EXTS

        if not (is_table or is_document):
            return JsonResponse(
                {"status": "error",
                 "message": f"Extensão não suportada: {ext}. "
                            f"Tabelas: {sorted(TABLE_EXTS)}. "
                            f"Documentos: PDF/DOCX/PPT/imagens."},
                status=400,
            )

        if is_document:
            return _handle_document_upload(
                f, conv_id=conv_id, agent_slug=agent_slug, user_note=user_note,
            )

        try:
            df = _read_table(f, f.name)
        except Exception as e:
            return JsonResponse(
                {"status": "error", "message": f"Falha ao ler o arquivo: {e}"},
                status=400,
            )

        if len(df) > UPLOAD_MAX_ROWS:
            df = df.head(UPLOAD_MAX_ROWS)
            truncated = True
        else:
            truncated = False

        # Recupera/cria a conversa
        conv: Conversation | None = None
        if conv_id:
            conv = Conversation.objects.filter(id=conv_id).first()
        if conv is None:
            agent = None
            if agent_slug:
                agent = Agent.objects.filter(slug=agent_slug).first()
            if agent is None:
                agent = Agent.objects.filter(is_default=True).first() or Agent.objects.first()
            if agent is None:
                return JsonResponse(
                    {"status": "error", "message": "Nenhum agente cadastrado."},
                    status=400,
                )
            title = f"📎 {f.name}"[:120]
            conv = Conversation.objects.create(title=title, agent=agent)

        # Salva o dataset corrente na sessão (mesma convenção das tools de análise)
        state = dict(conv.state or {})

        # Preserva dataset anterior como nomeado (multi-dataset)
        prev_source = state.get("athena_last_source") or {}
        prev_records = state.get("athena_last_result")
        if prev_records and prev_source:
            prev_name = prev_source.get("filename") or prev_source.get("query", "")[:40] or "anterior"
            prev_name = prev_name.replace(" ", "_").replace("/", "_")[:60]
            if "named_datasets" not in state:
                state["named_datasets"] = {}
            state["named_datasets"][prev_name] = prev_records

        state["athena_last_result"] = _df_to_records_jsonsafe(df)
        state["athena_last_columns"] = list(df.columns)
        state["athena_last_source"] = {"kind": "upload", "filename": f.name}
        if os.path.splitext(f.name)[1].lower() in {".xlsx", ".xls"}:
            sheet_name = df.attrs.get("sheet_name") if hasattr(df, "attrs") else None
            if sheet_name:
                state["athena_last_source"]["sheet_name"] = sheet_name
            # Persiste todas as abas como datasets nomeados. O dataset atual
            # continua sendo a aba escolhida pelo upload para não alterar a UI,
            # mas o agente pode mudar para qualquer outra aba via pandas.
            sheet_datasets = df.attrs.get("sheet_datasets", {}) if hasattr(df, "attrs") else {}
            if sheet_datasets:
                named = state.setdefault("named_datasets", {})
                workbook = {"sheets": [], "datasets": {}}
                for name, sheet_df in sheet_datasets.items():
                    if len(sheet_df) > UPLOAD_MAX_ROWS:
                        sheet_df = sheet_df.head(UPLOAD_MAX_ROWS)
                    dataset_name = f"{f.name}::{name}"
                    named[dataset_name] = _df_to_records_jsonsafe(sheet_df)
                    workbook["sheets"].append(str(name))
                    workbook["datasets"][str(name)] = dataset_name
                workbook["active_sheet"] = sheet_name
                state.setdefault("excel_workbooks", {})[f.name] = workbook
        conv.state = state
        conv.save()

        # Conteúdo "para o modelo" — curto, sem o dataset todo
        summary = _llm_summary_for_table(f.name, df)
        if user_note:
            summary = f"{user_note}\n\n{summary}"
        if truncated:
            summary += f"\n\n_(Dataset truncado nas primeiras {UPLOAD_MAX_ROWS:,} linhas.)_"

        # Attachment "para a UI" — preview maior, mas ainda paginável
        preview = _df_to_records_jsonsafe(df.head(PREVIEW_ROWS_FOR_UI))
        attachment = {
            "kind": "table",
            "filename": f.name,
            "rows": int(len(df)),
            "columns": list(df.columns),
            "dtypes": {c: str(df[c].dtype) for c in df.columns},
            "preview": preview,
            "preview_rows": len(preview),
            "truncated": truncated,
        }

        msg = Message.objects.create(
            conversation=conv,
            role="user",
            content=summary,
            attachment=attachment,
        )

        return JsonResponse({
            "status": "success",
            "conversation_id": conv.id,
            "conversation_title": conv.title,
            "agent_slug": conv.agent.slug if conv.agent else None,
            "message": _message_payload(msg),
        })
    except Exception as e:
        import traceback
        return JsonResponse(
            {
                "status": "error",
                "message": f"{type(e).__name__}: {e}",
                "traceback": traceback.format_exc(),
            },
            status=500,
        )


@require_GET
def conversation_dataset(request, conv_id):
    """Pagina o dataset corrente da conversa para a UI (sem passar pela LLM)."""
    conv = get_object_or_404(Conversation, id=conv_id)
    rows = (conv.state or {}).get("athena_last_result") or []
    columns = (conv.state or {}).get("athena_last_columns") or []
    try:
        offset = max(0, int(request.GET.get("offset", 0)))
    except ValueError:
        offset = 0
    try:
        limit = max(1, min(500, int(request.GET.get("limit", 100))))
    except ValueError:
        limit = 100

    return JsonResponse({
        "total": len(rows),
        "offset": offset,
        "limit": limit,
        "columns": columns,
        "rows": rows[offset: offset + limit],
    })


# ── Download de exports ──────────────────────────────────────────────

_EXPORT_FILENAME_RE = re.compile(r"^[A-Za-z0-9_-]+\.(csv|xlsx|pdf|html|htm)$")


_EXPORT_CONTENT_TYPES = {
    "csv": "text/csv; charset=utf-8",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pdf": "application/pdf",
    "html": "text/html; charset=utf-8",
    "htm": "text/html; charset=utf-8",
}


@require_GET
def export_download(request, filename: str):
    """Serve arquivos da pasta exports/. Filename validado para evitar
    path traversal — apenas [a-zA-Z0-9_-].(csv|xlsx)."""
    if not _EXPORT_FILENAME_RE.match(filename):
        raise Http404("Nome de arquivo inválido")

    path = Path(settings.BASE_DIR) / "exports" / filename
    if not path.exists() or not path.is_file():
        raise Http404("Arquivo não encontrado")

    ext = path.suffix.lstrip(".").lower()
    content_type = _EXPORT_CONTENT_TYPES.get(ext, "application/octet-stream")

    response = FileResponse(
        path.open("rb"),
        as_attachment=True,
        filename=filename,
        content_type=content_type,
    )
    # Garante que o browser baixe (sem essa header alguns clients abrem inline).
    response["Content-Disposition"] = f'attachment; filename="{filename}"'
    return response


@require_GET
def conversation_artifact_download(request, conv_id: int, filename: str):
    """Serve uma saída da conversa, com fallback para a pasta legada."""
    if not _EXPORT_FILENAME_RE.match(filename):
        raise Http404("Nome de arquivo inválido")

    workspace = Path(settings.BASE_DIR) / "runtime" / "codex_sessions" / str(conv_id)
    output_dir = (workspace / "saida").resolve()
    legacy_dir = (workspace / "artefatos").resolve()
    path = (output_dir / filename).resolve()
    if path.parent != output_dir or not path.is_file():
        path = (legacy_dir / filename).resolve()
    if path.parent not in {output_dir, legacy_dir} or not path.is_file():
        raise Http404("Arquivo não encontrado")

    ext = path.suffix.lstrip(".").lower()
    response = FileResponse(
        path.open("rb"),
        as_attachment=True,
        filename=filename,
        content_type=_EXPORT_CONTENT_TYPES.get(ext, "application/octet-stream"),
    )
    response["Content-Disposition"] = f'attachment; filename="{filename}"'
    return response
